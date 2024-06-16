from pptx import Presentation
from pptx.util import Inches

# Criação da apresentação
prs = Presentation()

# Adicionando um slide
slide_layout = prs.slide_layouts[5]  # layout de título e conteúdo
slide = prs.slides.add_slide(slide_layout)

# Título do slide
title = slide.shapes.title
title.text = "AutoChoice: Sua Plataforma Simples para Escolher o Carro Ideal"

# Adicionando o conteúdo
content = slide.placeholders[1]
content.text = (
    "Objetivo:\n"
    "Facilitar a escolha de veículos, proporcionando uma busca intuitiva e informações detalhadas em linguagem acessível.\n\n"
    "Funcionalidades:\n"
    "- Busca Simples e Rápida: Encontre carros por marca, modelo, preço, e outros critérios.\n"
    "- Informações Detalhadas: Descrições completas com especificações técnicas e avaliações.\n"
    "- Comparação de Veículos: Compare até três veículos lado a lado.\n"
    "- Recomendações Personalizadas: Sugestões baseadas nas suas preferências.\n\n"
    "Público-Alvo:\n"
    "Indicado para quem não tem conhecimento técnico sobre automóveis, famílias e qualquer pessoa buscando um carro de forma descomplicada.\n\n"
    "Contato:\n"
    "www.autochoice.com | contato@autochoice.com"
)

# Salvando a apresentação
pptx_path = "/mnt/data/AutoChoice_Slide.pptx"
prs.save(pptx_path)

# Compactando o arquivo
import zipfile

zip_path = "/mnt/data/AutoChoice_Slide.zip"
with zipfile.ZipFile(zip_path, 'w') as zipf:
    zipf.write(pptx_path, arcname="AutoChoice_Slide.pptx")